"""
Generates the PPTX presentation for the AI course assignment.
Run once, then delete this script.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── colour palette ────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy   – slide background
C_ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy – content bg
C_GREEN     = RGBColor(0x0F, 0x9B, 0x58)   # green
C_ORANGE    = RGBColor(0xE6, 0x7E, 0x22)   # orange
C_RED       = RGBColor(0xC0, 0x39, 0x2B)   # red
C_BLUE      = RGBColor(0x2E, 0x86, 0xC1)   # blue
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LTGREY    = RGBColor(0xCC, 0xCC, 0xCC)
C_YELLOW    = RGBColor(0xF1, 0xC4, 0x0F)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── utilities ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=C_BG, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, x, y, w, h,
             size=20, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, x, y, w, h,
                  size=16, color=C_WHITE, bold_first=False, line_spacing=None):
    """lines = list of (text, bold, color_override)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bld, col = item, False, color
        else:
            text = item[0]
            bld  = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else color

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bld or (bold_first and first)
        run.font.color.rgb = col
    return txb

def bg(slide):
    """Fill slide background with dark navy."""
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.2, C_ACCENT)
    add_text(slide, title, 0.3, 0.08, 12, 0.6,
             size=32, bold=True, color=C_GREEN)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.68, 12, 0.42,
                 size=16, color=C_LTGREY)

def slide_number(slide, n):
    add_text(slide, str(n), 12.8, 7.1, 0.45, 0.3,
             size=12, color=C_LTGREY, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)

# big accent bar
add_rect(s, 0, 2.5, 13.33, 2.7, C_ACCENT)
add_rect(s, 0, 2.5, 0.12, 2.7, C_GREEN)   # left accent stripe

add_text(s, "Paieška būsenų erdvėje", 0.5, 2.65, 12.5, 1.1,
         size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "State-Space Search · AIMA architektūra", 0.5, 3.7, 12.5, 0.6,
         size=22, color=C_GREEN, align=PP_ALIGN.CENTER)
add_text(s, "EightPuzzle  ·  NQueensProblem  ·  WaterJugProblem", 0.5, 4.25, 12.5, 0.5,
         size=16, color=C_LTGREY, align=PP_ALIGN.CENTER)

add_text(s, "Dirbtinis Intelektas · 2025–2026", 0.5, 6.8, 12.5, 0.5,
         size=13, color=C_LTGREY, align=PP_ALIGN.CENTER)
slide_number(s, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – Turinys
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Turinys", "Kas bus pristatyta")

items = [
    ("1.  AIMA architektūra – Problem, Node, paieška",   C_BLUE),
    ("2.  EightPuzzle – aštuonių dėlionė",               C_GREEN),
    ("3.  NQueensProblem – N karalienių uždavinys",      C_ORANGE),
    ("4.  Paieška į plotį vs į gylį vs informuota",      C_YELLOW),
    ("5.  Euristikos",                                    C_LTGREY),
    ("6.  WaterJugProblem – savas uždavinys (1.4)",       C_RED),
    ("7.  Išvados",                                       C_LTGREY),
]
for i, (txt, col) in enumerate(items):
    add_rect(s, 0.6, 1.4 + i*0.76, 11.8, 0.62, C_ACCENT)
    add_rect(s, 0.6, 1.4 + i*0.76, 0.08, 0.62, col)
    add_text(s, txt, 0.85, 1.43 + i*0.76, 11.2, 0.55,
             size=18, color=C_WHITE)
slide_number(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – AIMA architektūra
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "AIMA architektūra", "search.py pagrindiniai komponentai")

# Left column – class boxes
boxes = [
    ("Problem  (abstrakti klasė)", C_GREEN,
     "initial, goal\nactions(state)\nresult(state, action)\ngoal_test(state)\nh(node)"),
    ("Node  (paieškos medžio mazgas)", C_BLUE,
     "state, parent, action\npath_cost, depth\nexpand(problem)\nsolution() · path()"),
    ("Paieškos algoritmai", C_ORANGE,
     "breadth_first_graph_search  → FIFO\ndepth_first_graph_search     → Stack\nbest_first_graph_search      → Priority Q"),
]
for i, (title, col, body) in enumerate(boxes):
    add_rect(s, 0.4, 1.35 + i*1.9, 6.2, 1.75, C_ACCENT)
    add_rect(s, 0.4, 1.35 + i*1.9, 0.1, 1.75, col)
    add_text(s, title, 0.6, 1.38 + i*1.9, 5.9, 0.45,
             size=16, bold=True, color=col)
    add_text(s, body, 0.6, 1.8 + i*1.9, 5.9, 1.1,
             size=13, color=C_LTGREY)

# Right column – flow diagram
add_text(s, "Agento vykdymo ciklas:", 7.0, 1.3, 5.8, 0.4,
         size=16, bold=True, color=C_WHITE)
steps = [
    ("1. Pradinė būsena  →  Node(initial)", C_GREEN),
    ("2. Dedi į frontier (eilė / stack)", C_BLUE),
    ("3. Išimi mazgą, tikrini goal_test()", C_YELLOW),
    ("4. expand(): generuoji vaikus", C_ORANGE),
    ("5. Vaikai → frontier", C_BLUE),
    ("6. Radus tikslą: solution()  ←  path()", C_GREEN),
]
for i, (txt, col) in enumerate(steps):
    add_rect(s, 7.0, 1.75 + i*0.82, 5.9, 0.65, C_ACCENT)
    add_rect(s, 7.0, 1.75 + i*0.82, 0.08, 0.65, col)
    add_text(s, txt, 7.15, 1.78 + i*0.82, 5.65, 0.58,
             size=13, color=C_WHITE)
slide_number(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – EightPuzzle: aprašymas + būsena
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "EightPuzzle – Aštuonių dėlionė", "Uždavinio aprašymas ir būsenos reprezentacija")

# Left: description
add_rect(s, 0.35, 1.3, 6.3, 5.9, C_ACCENT)
add_text(s, "Uždavinys", 0.5, 1.35, 6.0, 0.45, size=17, bold=True, color=C_GREEN)
desc = (
    "3×3 lenta su 8 numeruotomis plytelėmis\n"
    "ir viena tuščia vieta (0).\n\n"
    "Tikslas: perstumdinėjant pasiekti\n"
    "tikslinę išdėstymą.\n\n"
    "Pradinė būsena:\n"
    "  (2, 4, 3, 1, 5, 6, 7, 8, 0)\n\n"
    "Tikslinė būsena:\n"
    "  (1, 2, 3, 4, 5, 6, 7, 8, 0)\n\n"
    "Veiksmai: UP · DOWN · LEFT · RIGHT\n"
    "(juda tuščia vieta)"
)
add_text(s, desc, 0.5, 1.82, 6.0, 5.3, size=15, color=C_WHITE)

# Right: state grid visualisation
add_rect(s, 6.9, 1.3, 6.1, 5.9, C_ACCENT)
add_text(s, "Būsenos vizualizacija", 7.05, 1.35, 5.8, 0.45,
         size=17, bold=True, color=C_ORANGE)

boards = [
    ("Pradinė  (0 žingsnis)",  [(2,4,3),(1,5,6),(7,8,0)], C_GREEN),
    ("Po LEFT  (1 žingsnis)",  [(2,4,3),(1,5,6),(7,0,8)], C_YELLOW),
    ("Po UP    (2 žingsnis)",  [(2,4,3),(1,0,6),(7,5,8)], C_ORANGE),
]
cell = 0.48
for bi, (label, board, col) in enumerate(boards):
    bx = 7.1 + bi * 2.0
    by = 1.85
    add_text(s, label, bx, by, 1.85, 0.3, size=10, color=col)
    for r in range(3):
        for c in range(3):
            val = board[r][c]
            tile_col = C_GREEN if val == 0 else C_ACCENT
            add_rect(s, bx + c*cell, by + 0.35 + r*cell, cell-0.03, cell-0.03, tile_col)
            txt = "_" if val == 0 else str(val)
            tcol = C_GREEN if val == 0 else C_WHITE
            add_text(s, txt, bx + c*cell + 0.02, by + 0.35 + r*cell + 0.02,
                     cell-0.07, cell-0.07, size=17, bold=True, color=tcol,
                     align=PP_ALIGN.CENTER)

# State tuple table
add_text(s, "Kortežo reikšmės:", 7.05, 4.05, 5.8, 0.35, size=14, bold=True, color=C_WHITE)
table_lines = [
    "Žingsnis  Veiksmas     Kortežas",
    "0  (pradinė)  —          (2,4,3,1,5,6,7,8,0)",
    "1            LEFT       (2,4,3,1,5,6,7,0,8)",
    "2            UP         (2,4,3,1,0,6,7,5,8)",
    "3            LEFT       (2,4,3,0,1,6,7,5,8)",
    "4            UP         (0,4,3,2,1,6,7,5,8)",
]
for i, line in enumerate(table_lines):
    col = C_YELLOW if i == 0 else C_LTGREY
    sz  = 12 if i == 0 else 11
    add_text(s, line, 7.05, 4.42 + i*0.37, 5.8, 0.34, size=sz, color=col)
slide_number(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – NQueensProblem: aprašymas + būsena
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "NQueensProblem – N karalienių uždavinys", "Būsenų reprezentacija ir veiksmai")

add_rect(s, 0.35, 1.3, 6.3, 5.9, C_ACCENT)
add_text(s, "Uždavinys", 0.5, 1.35, 6.0, 0.45, size=17, bold=True, color=C_GREEN)
desc = (
    "N karalienių išdėstymas N×N šachmatų\n"
    "lentoje taip, kad nė viena\n"
    "nekerstu kitos.\n\n"
    "N = 4  →  lenta 4×4\n\n"
    "Pradinė būsena:\n"
    "  (-1, -1, -1, -1)\n"
    "  (visi stulpeliai tušti)\n\n"
    "Tikslinė būsena:\n"
    "  visi stulpeliai užpildyti\n"
    "  ir nėra konfliktų\n\n"
    "Veiksmas: eilės numeris (0–N-1)\n"
    "į kurią dedama karalienė"
)
add_text(s, desc, 0.5, 1.82, 6.0, 5.3, size=15, color=C_WHITE)

# Right: board states
add_rect(s, 6.9, 1.3, 6.1, 5.9, C_ACCENT)
add_text(s, "Būsenų eiga (N=4)", 7.05, 1.35, 5.8, 0.45,
         size=17, bold=True, color=C_ORANGE)

states_4q = [
    ("(-1,-1,-1,-1)  pradinė",   [-1,-1,-1,-1]),
    ("(1,-1,-1,-1)   1 veiksmas",[ 1,-1,-1,-1]),
    ("(1,3,-1,-1)    2 veiksmas",[ 1, 3,-1,-1]),
    ("(1,3,0,-1)     3 veiksmas",[ 1, 3, 0,-1]),
    ("(1,3,0,2)      tikslinė",  [ 1, 3, 0, 2]),
]
cell = 0.38
for si, (label, state) in enumerate(states_4q):
    bx = 7.1
    by = 1.88 + si * 1.1
    col = C_GREEN if si == 0 else (C_RED if si == 4 else C_YELLOW)
    add_text(s, label, bx, by, 5.5, 0.28, size=10, color=col)
    for c in range(4):
        for r in range(4):
            has_q = (state[c] == r) and state[c] != -1
            tile_col = C_GREEN if has_q else (RGBColor(0x2C,0x3E,0x50) if (r+c)%2==0 else RGBColor(0x1A,0x27,0x38))
            add_rect(s, bx + c*cell, by + 0.3 + r*cell, cell-0.02, cell-0.02, tile_col)
            if has_q:
                add_text(s, "Q", bx + c*cell+0.01, by + 0.3 + r*cell,
                         cell-0.02, cell-0.02, size=11, bold=True,
                         color=C_WHITE, align=PP_ALIGN.CENTER)
slide_number(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – BFS vs DFS vs Best-First
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Paieškos algoritmai", "BFS · DFS · Best-First – kaip eilė lemia strategiją")

algos = [
    ("Paieška į plotį\n(BFS)", C_BLUE,
     "Eilė: FIFO  (deque)",
     "• Išimamas PIRMASIS elementas (popleft)\n"
     "• Pirma aplankomi VISI to gylio mazgai\n"
     "• Garantuoja TRUMPIAUSIĄ kelią\n"
     "• Naudoja daug atminties"),
    ("Paieška į gylį\n(DFS)", C_ORANGE,
     "Eilė: LIFO Stack  (list)",
     "• Išimamas PASKUTINIS elementas (pop)\n"
     "• Eina KIEK GALIMA giliau\n"
     "• Neranda trumpiausio kelio\n"
     "• Naudoja mažai atminties"),
    ("Informuota paieška\n(Best-First)", C_GREEN,
     "Eilė: Prioritetinė (heapq)",
     "• Išimamas MAŽIAUSIOS h(n) mazgas\n"
     "• Naudoja heuristinę funkciją\n"
     "• Greitesnė nei BFS didelėse erdvėse\n"
     "• Nėra garantuotas trumpiausias kelias"),
]

for i, (title, col, queue_txt, bullets) in enumerate(algos):
    x = 0.35 + i * 4.35
    add_rect(s, x, 1.3, 4.15, 5.9, C_ACCENT)
    add_rect(s, x, 1.3, 0.1, 5.9, col)
    add_text(s, title, x+0.15, 1.35, 3.9, 0.75, size=16, bold=True, color=col)
    add_rect(s, x+0.15, 2.1, 3.85, 0.45, RGBColor(0x0D,0x0D,0x1A))
    add_text(s, queue_txt, x+0.2, 2.13, 3.75, 0.38, size=13, color=C_YELLOW)
    add_text(s, bullets, x+0.15, 2.65, 3.9, 4.3, size=14, color=C_WHITE)

# Bottom code snippet
add_rect(s, 0.35, 7.0, 12.7, 0.38, RGBColor(0x0D,0x0D,0x1A))
add_text(s,
    "frontier = deque([Node(initial)])  # BFS     |    frontier = [Node(initial)]  # DFS",
    0.45, 7.02, 12.5, 0.34, size=11, color=C_YELLOW)
slide_number(s, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – Euristikos
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Euristikos", "h(n) – kiek liko iki tikslo?")

# Left
add_rect(s, 0.35, 1.3, 6.0, 5.9, C_ACCENT)
add_text(s, "EightPuzzle heuristikos", 0.5, 1.35, 5.7, 0.45,
         size=17, bold=True, color=C_GREEN)

ep_txt = (
    "1. Neteisingų plytelių skaičius\n"
    "   (Misplaced Tiles):\n\n"
    "   h(n) = kiek plytelių ne savo vietoje\n\n"
    "   Pavyzdys:\n"
    "   Esama:   (2,4,3,1,5,6,7,8,0)\n"
    "   Tikslas: (1,2,3,4,5,6,7,8,0)\n"
    "   2≠1, 4≠2, 1≠4  →  h = 3\n\n"
    "2. Manheteno atstumas\n"
    "   (Manhattan Distance):\n\n"
    "   h(n) = Σ |eilutė_d - eilutė_t|\n"
    "            + |stulpelis_d - stulpelis_t|\n\n"
    "   Tikslesnė, bet sudėtingesnė."
)
add_text(s, ep_txt, 0.5, 1.85, 5.7, 5.3, size=13, color=C_WHITE)

# Right
add_rect(s, 6.65, 1.3, 6.3, 5.9, C_ACCENT)
add_text(s, "NQueensProblem heuristika", 6.8, 1.35, 6.0, 0.45,
         size=17, bold=True, color=C_ORANGE)

nq_txt = (
    "Kertančių porų skaičius:\n\n"
    "   h(n) = kiek karalienių porų\n"
    "           kertasi tarpusavyje\n\n"
    "   Kuo mažiau — tuo geriau\n"
    "   h = 0  →  tikslinė būsena!\n\n"
    "Konfliktas egzistuoja, jei:\n"
    "  • ta pati eilutė\n"
    "  • tas pats stulpelis\n"
    "  • ta pati įstrižainė\n"
    "    (|r1-r2| == |c1-c2|)\n\n"
    "WaterJugProblem heuristika:\n\n"
    "   h(n) = |a - 2|\n"
    "   (kiek litrų trūksta iki 2L)"
)
add_text(s, nq_txt, 6.8, 1.85, 6.0, 5.3, size=13, color=C_WHITE)
slide_number(s, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – WaterJugProblem aprašymas
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "WaterJugProblem – Savas uždavinys (1.4)",
           "Vandens kibirų problema · AIMA architektūra")

# AI prompt box
add_rect(s, 0.35, 1.32, 12.65, 1.0, RGBColor(0x0D, 0x0D, 0x1A))
add_rect(s, 0.35, 1.32, 0.1, 1.0, C_YELLOW)
add_text(s, "AI Prompt (GitHub Copilot / Claude):", 0.55, 1.35, 12.2, 0.32,
         size=12, bold=True, color=C_YELLOW)
add_text(s,
    '"Using the AIMA search.py architecture (Problem base class), implement the Water Jug Problem: '
    'Jug A=4L, Jug B=3L, goal=2L in A. '
    'Solve with BFS, DFS, best-first. Print step-by-step, draw networkx state-space graph."',
    0.55, 1.67, 12.2, 0.58, size=11, color=C_LTGREY, italic=True)

# Left: problem description
add_rect(s, 0.35, 2.45, 6.0, 4.8, C_ACCENT)
add_text(s, "Uždavinio aprašymas", 0.5, 2.5, 5.7, 0.4, size=16, bold=True, color=C_GREEN)
wd_txt = (
    "Kibiras A: talpa 4 litrai\n"
    "Kibiras B: talpa 3 litrai\n\n"
    "Pradinė būsena: (0, 0)\n"
    "  A=0L, B=0L – abu tušti\n\n"
    "Tikslinė būsena: A == 2L\n\n"
    "Veiksmai (6):\n"
    "  Fill A    – pripilti A iki 4L\n"
    "  Fill B    – pripilti B iki 3L\n"
    "  Empty A   – išpilti A\n"
    "  Empty B   – išpilti B\n"
    "  Pour A→B  – pilti iš A į B\n"
    "  Pour B→A  – pilti iš B į A"
)
add_text(s, wd_txt, 0.5, 2.95, 5.7, 4.2, size=14, color=C_WHITE)

# Right: step-by-step solution
add_rect(s, 6.55, 2.45, 6.45, 4.8, C_ACCENT)
add_text(s, "BFS sprendimas (6 žingsniai)", 6.7, 2.5, 6.1, 0.4,
         size=16, bold=True, color=C_ORANGE)
steps = [
    ("Pradžia",     "(0, 0)",  "A=0L  B=0L"),
    ("Fill A",      "(4, 0)",  "A=4L  B=0L"),
    ("Pour A→B",    "(1, 3)",  "A=1L  B=3L"),
    ("Empty B",     "(1, 0)",  "A=1L  B=0L"),
    ("Pour A→B",    "(0, 1)",  "A=0L  B=1L"),
    ("Fill A",      "(4, 1)",  "A=4L  B=1L"),
    ("Pour A→B",    "(2, 3)",  "A=2L  B=3L  ✓ TIKSLAS"),
]
for i, (action, state, label) in enumerate(steps):
    col = C_RED if i == 6 else (C_GREEN if i == 0 else C_LTGREY)
    y = 2.98 + i*0.6
    add_rect(s, 6.7, y, 6.1, 0.53, RGBColor(0x0D,0x0D,0x1A) if i in (0,6) else C_ACCENT)
    add_text(s, f"  {i}.  {action:<12} {state:<10}  {label}", 6.75, y+0.04,
             5.9, 0.46, size=12, color=col)
slide_number(s, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – WaterJugProblem kodas
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "WaterJugProblem – Kodo architektūra",
           "Problem klasės implementacija · AIMA šablonas")

code_blocks = [
    ("__init__  –  pradinė būsena", C_GREEN,
     "def __init__(self):\n"
     "    self.CAP_A = 4    # Kibiras A: 4L\n"
     "    self.CAP_B = 3    # Kibiras B: 3L\n"
     "    super().__init__(initial=(0, 0))"),

    ("actions(state)  –  galimi veiksmai", C_BLUE,
     "def actions(self, state):\n"
     "    a, b = state\n"
     "    if a < CAP_A:  possible.append('Fill A')\n"
     "    if b < CAP_B:  possible.append('Fill B')\n"
     "    if a > 0:      possible.append('Empty A')\n"
     "    if a > 0 and b < CAP_B: possible.append('Pour A->B')\n"
     "    ..."),

    ("result(state, action)  –  nauja būsena", C_ORANGE,
     "def result(self, state, action):\n"
     "    if action == 'Fill A':    return (CAP_A, b)\n"
     "    if action == 'Pour A->B':\n"
     "        pour = min(a, CAP_B - b)\n"
     "        return (a - pour, b + pour)"),

    ("goal_test  +  h  –  tikslas ir heuristika", C_YELLOW,
     "def goal_test(self, state):\n"
     "    return state[0] == 2   # A == 2L\n\n"
     "def h(self, node):\n"
     "    return abs(node.state[0] - 2)"),
]

for i, (title, col, code) in enumerate(code_blocks):
    row, cidx = divmod(i, 2)
    x = 0.35 + cidx * 6.55
    y = 1.35 + row * 3.0
    add_rect(s, x, y, 6.3, 2.85, C_ACCENT)
    add_rect(s, x, y, 0.08, 2.85, col)
    add_text(s, title, x+0.15, y+0.05, 6.0, 0.38, size=13, bold=True, color=col)
    add_rect(s, x+0.12, y+0.48, 6.03, 2.28, RGBColor(0x0D,0x0D,0x1A))
    add_text(s, code, x+0.18, y+0.52, 5.9, 2.2, size=11, color=C_LTGREY)
slide_number(s, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – Grafai
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Būsenų erdvės grafai", "Vizualizacija su networkx + matplotlib")

imgs = [
    ("eight_puzzle_graph.png",  "EightPuzzle – BFS grafas",  C_GREEN,   0.35, 1.35, 6.1, 5.5),
    ("nqueens_graph.png",       "NQueens – Best-First grafas", C_ORANGE, 6.75, 1.35, 6.1, 5.5),
]
for fname, caption, col, x, y, w, h_i in imgs:
    add_rect(s, x, y, w, h_i, C_ACCENT)
    add_text(s, caption, x+0.1, y+0.05, w-0.2, 0.38, size=14, bold=True, color=col)
    if os.path.exists(fname):
        try:
            slide = s
            pic = slide.shapes.add_picture(
                fname,
                Inches(x+0.08), Inches(y+0.48),
                Inches(w-0.16), Inches(h_i-0.58)
            )
        except Exception:
            add_text(s, f"[{fname}]", x+0.2, y+0.55, w-0.3, h_i-0.65,
                     size=12, color=C_LTGREY, align=PP_ALIGN.CENTER)
    else:
        add_text(s, f"[{fname} – paleiskite skriptą]",
                 x+0.2, y+2.5, w-0.3, 0.5, size=11, color=C_LTGREY, align=PP_ALIGN.CENTER)

slide_number(s, 10)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – WaterJug grafas
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "WaterJugProblem – Būsenų erdvės grafas",
           "14 būsenų · 13 perėjimų · BFS sprendimas 6 žingsniais")

fname = "water_jug_graph.png"
add_rect(s, 0.35, 1.3, 12.65, 5.95, C_ACCENT)
if os.path.exists(fname):
    try:
        s.shapes.add_picture(fname, Inches(0.4), Inches(1.35), Inches(12.55), Inches(5.85))
    except Exception:
        add_text(s, "[water_jug_graph.png]", 0.5, 4.0, 12.3, 0.6,
                 size=14, color=C_LTGREY, align=PP_ALIGN.CENTER)
else:
    add_text(s, "[Paleiskite WaterJugProblem.py kad sugeneruotumėte grafą]",
             0.5, 4.0, 12.3, 0.6, size=14, color=C_LTGREY, align=PP_ALIGN.CENTER)
slide_number(s, 11)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 – Išvados
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Išvados", "Ko išmokome")

conclusions = [
    (C_GREEN,  "AIMA architektūra",
     "Problem → Node → paieškos algoritmas.\n"
     "Kiekvienas uždavinys implementuoja actions(), result(), goal_test(), h()."),
    (C_BLUE,   "Eilė lemia strategiją",
     "FIFO deque = BFS (plotis) · LIFO stack = DFS (gylis)\n"
     "Prioritetinė eilė + h(n) = informuota paieška."),
    (C_ORANGE, "Euristikos pagreitina paiešką",
     "Misplaced tiles (EightPuzzle) · kertančios poros (NQueens)\n"
     "· |a-2| (WaterJug). Geresnis h → mažiau aplankytų mazgų."),
    (C_YELLOW, "WaterJugProblem (1.4)",
     "Savas uždavinys sukurtas su AI. 14 būsenų, sprendimas 6 žingsniais.\n"
     "Pilnai integruotas su AIMA search.py architektūra."),
]

for i, (col, title, body) in enumerate(conclusions):
    row, cidx = divmod(i, 2)
    x = 0.35 + cidx * 6.55
    y = 1.35 + row * 2.85
    add_rect(s, x, y, 6.3, 2.7, C_ACCENT)
    add_rect(s, x, y, 0.1, 2.7, col)
    add_text(s, title, x+0.18, y+0.08, 5.95, 0.42, size=17, bold=True, color=col)
    add_text(s, body,  x+0.18, y+0.55, 5.95, 2.0,  size=14, color=C_WHITE)
slide_number(s, 12)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = "Paieška_Busenų_Erdveje.pptx"
prs.save(out)
print(f"Saved → {out}")
